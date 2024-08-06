from pptx import Presentation
from pptx.util import Inches

# Create a PowerPoint presentation object
prs = Presentation()

# Title Slide
slide = prs.slides.add_slide(prs.slide_layouts[0])
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "Instagram Automation Tool"
subtitle.text = "Automating Repetitive Tasks Using Java and Selenium"

# Introduction Slide
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
title.text = "Introduction"
content = slide.placeholders[1].text_frame
content.text = (
    "This project aims to automate repetitive tasks on Instagram such as liking posts, "
    "commenting, and following users using Java and Selenium WebDriver. The automation "
    "helps in saving time and effort for digital marketers and social media enthusiasts."
)

# Team Slide
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
title.text = "Team Members"
content = slide.placeholders[1].text_frame
content.text = (
    "Anurag Sharma - Project Lead\n"
    "Kevin Paulose - Developer\n"
    "Kiran Suresh - Developer"
)

# Project Overview Slide
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
title.text = "Project Overview"
content = slide.placeholders[1].text_frame
content.text = (
    "The project automates the following tasks:\n"
    "1. Login to Instagram\n"
    "2. Navigate to explore and profile pages\n"
    "3. Automate liking, commenting, and following\n"
    "4. Store metadata for further analysis"
)

# Setup and Configuration Slide
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
title.text = "Setup and Configuration"
content = slide.placeholders[1].text_frame
content.text = (
    "1. Install Java and Maven\n"
    "2. Install Selenium WebDriver\n"
    "3. Download and configure Geckodriver\n"
    "4. Set up project directory and configuration files"
)

# Login Automation Slide
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
title.text = "Login Automation"
content = slide.placeholders[1].text_frame
content.text = (
    "1. Open Instagram login page\n"
    "2. Enter username and password\n"
    "3. Click login button\n"
    "4. Handle login verification if needed"
)

# Actions Automation Slide
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
title.text = "Actions Automation"
content = slide.placeholders[1].text_frame
content.text = (
    "1. Navigate to explore page\n"
    "2. Interact with posts:\n"
    "   - Like posts\n"
    "   - Comment on posts\n"
    "   - Follow users\n"
    "3. Store metadata for each interaction"
)

# Data Handling Slide
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
title.text = "Data Handling"
content = slide.placeholders[1].text_frame
content.text = (
    "1. Extract metadata from posts\n"
    "2. Store metadata in CSV files\n"
    "Metadata includes:\n"
    "   - Profile name\n"
    "   - Number of likes\n"
    "   - Comments\n"
    "   - Date posted\n"
    "   - Post URL"
)

# Challenges and Solutions Slide
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
title.text = "Challenges and Solutions"
content = slide.placeholders[1].text_frame
content.text = (
    "1. Handling dynamic content loading\n"
    "2. Managing login verification\n"
    "3. Avoiding detection by Instagram\n"
    "4. Ensuring data accuracy"
)

# Future Enhancements Slide
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
title.text = "Future Enhancements"
content = slide.placeholders[1].text_frame
content.text = (
    "1. Integrate machine learning for sentiment analysis\n"
    "2. Enhance data visualization\n"
    "3. Add support for other social media platforms\n"
    "4. Implement a user-friendly interface"
)

# Conclusion Slide
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
title.text = "Conclusion"
content = slide.placeholders[1].text_frame
content.text = (
    "The Instagram Automation Tool successfully automates repetitive tasks, saving time "
    "and increasing efficiency for users. The project demonstrates the power of Java and "
    "Selenium for web automation and sets the stage for future enhancements."
)

# Save the presentation
pptx_file = "Instagram_Automation_Tool_Presentation.pptx"
prs.save(pptx_file)
